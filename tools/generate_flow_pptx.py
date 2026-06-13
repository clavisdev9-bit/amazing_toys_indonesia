"""
Generate ATF2026_Flow_Charts.pptx
Amazing Toys Fair 2026 — Process Flow Diagrams
"""

import io
import tempfile
import os
from svglib.svglib import svg2rlg
from reportlab.graphics import renderPM
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── SVG DEFINITIONS ──────────────────────────────────────────────────────────

SVG_GLOBAL = """<?xml version="1.0" encoding="UTF-8"?>
<svg viewBox="0 0 820 680" xmlns="http://www.w3.org/2000/svg">
  <rect width="820" height="680" fill="white"/>
  <defs>
    <marker id="arr" markerWidth="8" markerHeight="6" refX="8" refY="3" orient="auto"><polygon points="0 0, 8 3, 0 6" fill="#6b7280"/></marker>
    <marker id="arr-red" markerWidth="8" markerHeight="6" refX="8" refY="3" orient="auto"><polygon points="0 0, 8 3, 0 6" fill="#ef4444"/></marker>
    <marker id="arr-green" markerWidth="8" markerHeight="6" refX="8" refY="3" orient="auto"><polygon points="0 0, 8 3, 0 6" fill="#16a34a"/></marker>
  </defs>
  <rect x="0" y="0" width="820" height="30" fill="#1e3a5f"/>
  <text x="410" y="20" text-anchor="middle" font-size="14" font-weight="bold" fill="white" font-family="Arial">GLOBAL ORDER FLOW — Amazing Toys Fair 2026</text>
  <!-- Swim lane backgrounds -->
  <rect x="0" y="30" width="820" height="120" fill="#dbeafe" opacity="0.5"/>
  <rect x="0" y="150" width="820" height="120" fill="#dcfce7" opacity="0.5"/>
  <rect x="0" y="270" width="820" height="120" fill="#fef9c3" opacity="0.5"/>
  <rect x="0" y="390" width="820" height="100" fill="#fee2e2" opacity="0.5"/>
  <rect x="0" y="490" width="820" height="100" fill="#ede9fe" opacity="0.5"/>
  <rect x="0" y="590" width="820" height="90" fill="#f3f4f6" opacity="0.5"/>
  <!-- Lane labels -->
  <text x="14" y="100" font-size="10" font-weight="bold" fill="#1e40af" font-family="Arial" transform="rotate(-90,14,100)">CUSTOMER</text>
  <text x="14" y="220" font-size="10" font-weight="bold" fill="#15803d" font-family="Arial" transform="rotate(-90,14,220)">HELPER</text>
  <text x="14" y="340" font-size="10" font-weight="bold" fill="#a16207" font-family="Arial" transform="rotate(-90,14,340)">CASHIER</text>
  <text x="14" y="448" font-size="10" font-weight="bold" fill="#b91c1c" font-family="Arial" transform="rotate(-90,14,448)">SYSTEM</text>
  <text x="14" y="552" font-size="10" font-weight="bold" fill="#6d28d9" font-family="Arial" transform="rotate(-90,14,552)">ADMIN</text>
  <text x="14" y="645" font-size="10" font-weight="bold" fill="#374151" font-family="Arial" transform="rotate(-90,14,645)">STATUS</text>
  <!-- CUSTOMER ROW -->
  <ellipse cx="80" cy="90" rx="38" ry="18" fill="#2563eb"/>
  <text x="80" y="94" text-anchor="middle" font-size="10" font-weight="bold" fill="white" font-family="Arial">Mulai</text>
  <rect x="140" y="72" width="110" height="36" rx="6" fill="#dbeafe" stroke="#2563eb" stroke-width="1.5"/>
  <text x="195" y="87" text-anchor="middle" font-size="10" fill="#1e40af" font-family="Arial">Registrasi /</text>
  <text x="195" y="100" text-anchor="middle" font-size="10" fill="#1e40af" font-family="Arial">Login</text>
  <rect x="280" y="72" width="90" height="36" rx="6" fill="#dbeafe" stroke="#2563eb" stroke-width="1.5"/>
  <text x="325" y="87" text-anchor="middle" font-size="10" fill="#1e40af" font-family="Arial">Browse &amp;</text>
  <text x="325" y="100" text-anchor="middle" font-size="10" fill="#1e40af" font-family="Arial">Pilih Produk</text>
  <rect x="400" y="72" width="90" height="36" rx="6" fill="#dbeafe" stroke="#2563eb" stroke-width="1.5"/>
  <text x="445" y="87" text-anchor="middle" font-size="10" fill="#1e40af" font-family="Arial">Checkout /</text>
  <text x="445" y="100" text-anchor="middle" font-size="10" fill="#1e40af" font-family="Arial">Submit Order</text>
  <polygon points="570,68 620,90 570,112 520,90" fill="#fde68a" stroke="#d97706" stroke-width="1.5"/>
  <text x="570" y="86" text-anchor="middle" font-size="9" font-weight="bold" fill="#92400e" font-family="Arial">Order</text>
  <text x="570" y="98" text-anchor="middle" font-size="9" font-weight="bold" fill="#92400e" font-family="Arial">Mode?</text>
  <rect x="680" y="72" width="110" height="36" rx="6" fill="#dbeafe" stroke="#2563eb" stroke-width="1.5"/>
  <text x="735" y="87" text-anchor="middle" font-size="10" fill="#1e40af" font-family="Arial">Terima QR &amp;</text>
  <text x="735" y="100" text-anchor="middle" font-size="10" fill="#1e40af" font-family="Arial">Bayar</text>
  <!-- Customer arrows -->
  <line x1="118" y1="90" x2="138" y2="90" stroke="#6b7280" stroke-width="1.5" marker-end="url(#arr)"/>
  <line x1="250" y1="90" x2="278" y2="90" stroke="#6b7280" stroke-width="1.5" marker-end="url(#arr)"/>
  <line x1="370" y1="90" x2="398" y2="90" stroke="#6b7280" stroke-width="1.5" marker-end="url(#arr)"/>
  <line x1="490" y1="90" x2="518" y2="90" stroke="#6b7280" stroke-width="1.5" marker-end="url(#arr)"/>
  <line x1="622" y1="90" x2="678" y2="90" stroke="#6b7280" stroke-width="1.5" marker-end="url(#arr)"/>
  <!-- HELPER ROW -->
  <rect x="140" y="162" width="110" height="36" rx="6" fill="#dcfce7" stroke="#16a34a" stroke-width="1.5"/>
  <text x="195" y="177" text-anchor="middle" font-size="10" fill="#15803d" font-family="Arial">Buat Order</text>
  <text x="195" y="190" text-anchor="middle" font-size="9" fill="#15803d" font-family="Arial">(HELPER_INPUT)</text>
  <rect x="400" y="162" width="110" height="36" rx="6" fill="#dcfce7" stroke="#16a34a" stroke-width="1.5"/>
  <text x="455" y="177" text-anchor="middle" font-size="10" fill="#15803d" font-family="Arial">Review Queue</text>
  <text x="455" y="190" text-anchor="middle" font-size="9" fill="#15803d" font-family="Arial">Approval</text>
  <polygon points="580,158 630,180 580,202 530,180" fill="#dcfce7" stroke="#16a34a" stroke-width="1.5"/>
  <text x="580" y="176" text-anchor="middle" font-size="9" font-weight="bold" fill="#15803d" font-family="Arial">Approve</text>
  <text x="580" y="188" text-anchor="middle" font-size="9" font-weight="bold" fill="#15803d" font-family="Arial">/ Reject?</text>
  <rect x="680" y="162" width="110" height="36" rx="6" fill="#dcfce7" stroke="#16a34a" stroke-width="1.5"/>
  <text x="735" y="177" text-anchor="middle" font-size="10" fill="#15803d" font-family="Arial">Konfirmasi</text>
  <text x="735" y="190" text-anchor="middle" font-size="10" fill="#15803d" font-family="Arial">Handover</text>
  <!-- Reject path -->
  <line x1="580" y1="202" x2="580" y2="238" stroke="#ef4444" stroke-width="1.5" marker-end="url(#arr-red)"/>
  <rect x="540" y="238" width="80" height="22" rx="6" fill="#fee2e2" stroke="#ef4444" stroke-width="1.2"/>
  <text x="580" y="253" text-anchor="middle" font-size="9" fill="#b91c1c" font-family="Arial">CANCELLED</text>
  <text x="592" y="208" font-size="8" fill="#ef4444" font-family="Arial">REJECT</text>
  <!-- Mode routing -->
  <line x1="570" y1="112" x2="570" y2="148" stroke="#16a34a" stroke-width="1.5" stroke-dasharray="4"/>
  <line x1="570" y1="148" x2="252" y2="160" stroke="#16a34a" stroke-width="1.5" marker-end="url(#arr-green)"/>
  <line x1="570" y1="112" x2="570" y2="150" stroke="#d97706" stroke-width="1.5" stroke-dasharray="4"/>
  <line x1="570" y1="150" x2="452" y2="160" stroke="#d97706" stroke-width="1.5" marker-end="url(#arr)"/>
  <text x="534" y="143" font-size="8" fill="#16a34a" font-family="Arial">HELPER_INPUT</text>
  <text x="534" y="130" font-size="8" fill="#d97706" font-family="Arial">APPROVE</text>
  <line x1="512" y1="180" x2="402" y2="180" stroke="#6b7280" stroke-width="1.5" marker-end="url(#arr)"/>
  <line x1="632" y1="170" x2="680" y2="90" stroke="#16a34a" stroke-width="1.3" stroke-dasharray="3" marker-end="url(#arr-green)"/>
  <!-- CASHIER ROW -->
  <rect x="400" y="282" width="110" height="36" rx="6" fill="#fef9c3" stroke="#ca8a04" stroke-width="1.5"/>
  <text x="455" y="297" text-anchor="middle" font-size="10" fill="#a16207" font-family="Arial">Scan QR /</text>
  <text x="455" y="310" text-anchor="middle" font-size="10" fill="#a16207" font-family="Arial">Konfirmasi</text>
  <rect x="540" y="282" width="110" height="36" rx="6" fill="#fef9c3" stroke="#ca8a04" stroke-width="1.5"/>
  <text x="595" y="297" text-anchor="middle" font-size="10" fill="#a16207" font-family="Arial">Proses</text>
  <text x="595" y="310" text-anchor="middle" font-size="10" fill="#a16207" font-family="Arial">Pembayaran</text>
  <rect x="680" y="282" width="110" height="36" rx="6" fill="#fef9c3" stroke="#ca8a04" stroke-width="1.5"/>
  <text x="735" y="297" text-anchor="middle" font-size="10" fill="#a16207" font-family="Arial">Update</text>
  <text x="735" y="310" text-anchor="middle" font-size="10" fill="#a16207" font-family="Arial">Status PAID</text>
  <line x1="510" y1="300" x2="538" y2="300" stroke="#6b7280" stroke-width="1.5" marker-end="url(#arr)"/>
  <line x1="650" y1="300" x2="678" y2="300" stroke="#6b7280" stroke-width="1.5" marker-end="url(#arr)"/>
  <line x1="735" y1="108" x2="735" y2="280" stroke="#6b7280" stroke-width="1.3" stroke-dasharray="3" marker-end="url(#arr)"/>
  <!-- SYSTEM ROW -->
  <rect x="140" y="402" width="120" height="36" rx="6" fill="#fee2e2" stroke="#ef4444" stroke-width="1.5"/>
  <text x="200" y="417" text-anchor="middle" font-size="10" fill="#b91c1c" font-family="Arial">Auto Expire</text>
  <text x="200" y="430" text-anchor="middle" font-size="9" fill="#b91c1c" font-family="Arial">TxnExpireJob</text>
  <rect x="400" y="402" width="120" height="36" rx="6" fill="#fee2e2" stroke="#ef4444" stroke-width="1.5"/>
  <text x="460" y="417" text-anchor="middle" font-size="10" fill="#b91c1c" font-family="Arial">WA Notif</text>
  <text x="460" y="430" text-anchor="middle" font-size="9" fill="#b91c1c" font-family="Arial">(WAHA)</text>
  <rect x="560" y="402" width="120" height="36" rx="6" fill="#fee2e2" stroke="#ef4444" stroke-width="1.5"/>
  <text x="620" y="417" text-anchor="middle" font-size="10" fill="#b91c1c" font-family="Arial">WebSocket</text>
  <text x="620" y="430" text-anchor="middle" font-size="9" fill="#b91c1c" font-family="Arial">Broadcast</text>
  <line x1="520" y1="420" x2="558" y2="420" stroke="#6b7280" stroke-width="1.5" marker-end="url(#arr)"/>
  <!-- ADMIN ROW -->
  <rect x="60" y="502" width="100" height="36" rx="6" fill="#ede9fe" stroke="#7c3aed" stroke-width="1.5"/>
  <text x="110" y="517" text-anchor="middle" font-size="10" fill="#6d28d9" font-family="Arial">Manage Users</text>
  <text x="110" y="530" text-anchor="middle" font-size="9" fill="#6d28d9" font-family="Arial">+ Config</text>
  <rect x="190" y="502" width="100" height="36" rx="6" fill="#ede9fe" stroke="#7c3aed" stroke-width="1.5"/>
  <text x="240" y="517" text-anchor="middle" font-size="10" fill="#6d28d9" font-family="Arial">Override</text>
  <text x="240" y="530" text-anchor="middle" font-size="9" fill="#6d28d9" font-family="Arial">Order</text>
  <rect x="320" y="502" width="100" height="36" rx="6" fill="#ede9fe" stroke="#7c3aed" stroke-width="1.5"/>
  <text x="370" y="517" text-anchor="middle" font-size="10" fill="#6d28d9" font-family="Arial">Sync Odoo</text>
  <text x="370" y="530" text-anchor="middle" font-size="9" fill="#6d28d9" font-family="Arial">+ Laporan</text>
  <!-- STATUS ROW -->
  <rect x="60" y="608" width="106" height="26" rx="13" fill="#e0e7ff" stroke="#6366f1" stroke-width="1.5"/>
  <text x="113" y="625" text-anchor="middle" font-size="9" font-weight="bold" fill="#4338ca" font-family="Arial">PENDING_APPROVAL</text>
  <rect x="190" y="608" width="80" height="26" rx="13" fill="#fef9c3" stroke="#ca8a04" stroke-width="1.5"/>
  <text x="230" y="625" text-anchor="middle" font-size="9" font-weight="bold" fill="#a16207" font-family="Arial">PENDING</text>
  <rect x="296" y="608" width="80" height="26" rx="13" fill="#dcfce7" stroke="#16a34a" stroke-width="1.5"/>
  <text x="336" y="625" text-anchor="middle" font-size="9" font-weight="bold" fill="#15803d" font-family="Arial">RESERVED</text>
  <rect x="402" y="608" width="80" height="26" rx="13" fill="#fee2e2" stroke="#ef4444" stroke-width="1.5"/>
  <text x="442" y="625" text-anchor="middle" font-size="9" font-weight="bold" fill="#b91c1c" font-family="Arial">EXPIRED</text>
  <rect x="508" y="608" width="60" height="26" rx="13" fill="#bbf7d0" stroke="#15803d" stroke-width="1.5"/>
  <text x="538" y="625" text-anchor="middle" font-size="9" font-weight="bold" fill="#14532d" font-family="Arial">PAID</text>
  <rect x="592" y="608" width="90" height="26" rx="13" fill="#fde68a" stroke="#d97706" stroke-width="1.5"/>
  <text x="637" y="625" text-anchor="middle" font-size="9" font-weight="bold" fill="#92400e" font-family="Arial">HANDED_OVER</text>
  <rect x="704" y="608" width="80" height="26" rx="13" fill="#bfdbfe" stroke="#2563eb" stroke-width="1.5"/>
  <text x="744" y="625" text-anchor="middle" font-size="9" font-weight="bold" fill="#1e40af" font-family="Arial">COMPLETED</text>
  <line x1="168" y1="621" x2="188" y2="621" stroke="#6b7280" stroke-width="1.2" marker-end="url(#arr)"/>
  <line x1="270" y1="621" x2="294" y2="621" stroke="#6b7280" stroke-width="1.2" marker-end="url(#arr)"/>
  <line x1="480" y1="621" x2="506" y2="621" stroke="#6b7280" stroke-width="1.2" marker-end="url(#arr)"/>
  <line x1="568" y1="621" x2="590" y2="621" stroke="#6b7280" stroke-width="1.2" marker-end="url(#arr)"/>
  <line x1="682" y1="621" x2="702" y2="621" stroke="#6b7280" stroke-width="1.2" marker-end="url(#arr)"/>
</svg>"""

SVG_HELPER_APPROVE = """<?xml version="1.0" encoding="UTF-8"?>
<svg viewBox="0 0 820 540" xmlns="http://www.w3.org/2000/svg">
  <rect width="820" height="540" fill="white"/>
  <defs>
    <marker id="a2" markerWidth="8" markerHeight="6" refX="8" refY="3" orient="auto"><polygon points="0 0,8 3,0 6" fill="#6b7280"/></marker>
    <marker id="a2r" markerWidth="8" markerHeight="6" refX="8" refY="3" orient="auto"><polygon points="0 0,8 3,0 6" fill="#ef4444"/></marker>
    <marker id="a2g" markerWidth="8" markerHeight="6" refX="8" refY="3" orient="auto"><polygon points="0 0,8 3,0 6" fill="#16a34a"/></marker>
  </defs>
  <rect x="0" y="0" width="820" height="26" fill="#6d28d9"/>
  <text x="410" y="18" text-anchor="middle" font-size="12" font-weight="bold" fill="white" font-family="Arial">MODE: HELPER_APPROVE (Model D) — Active Mode</text>
  <!-- Stage columns -->
  <rect x="8" y="32" width="130" height="500" rx="6" fill="#f5f3ff" stroke="#c4b5fd" stroke-width="1"/>
  <text x="73" y="52" text-anchor="middle" font-size="9" font-weight="bold" fill="#6d28d9" font-family="Arial">STAGE 1</text>
  <text x="73" y="64" text-anchor="middle" font-size="9" fill="#7c3aed" font-family="Arial">Customer Checkout</text>
  <rect x="148" y="32" width="130" height="500" rx="6" fill="#f0fdf4" stroke="#bbf7d0" stroke-width="1"/>
  <text x="213" y="52" text-anchor="middle" font-size="9" font-weight="bold" fill="#15803d" font-family="Arial">STAGE 2</text>
  <text x="213" y="64" text-anchor="middle" font-size="9" fill="#16a34a" font-family="Arial">Helper Approval</text>
  <rect x="288" y="32" width="130" height="500" rx="6" fill="#fffbeb" stroke="#fde68a" stroke-width="1"/>
  <text x="353" y="52" text-anchor="middle" font-size="9" font-weight="bold" fill="#a16207" font-family="Arial">STAGE 3</text>
  <text x="353" y="64" text-anchor="middle" font-size="9" fill="#ca8a04" font-family="Arial">QR &amp; Payment</text>
  <rect x="428" y="32" width="130" height="500" rx="6" fill="#fff7ed" stroke="#fed7aa" stroke-width="1"/>
  <text x="493" y="52" text-anchor="middle" font-size="9" font-weight="bold" fill="#c2410c" font-family="Arial">STAGE 4</text>
  <text x="493" y="64" text-anchor="middle" font-size="9" fill="#ea580c" font-family="Arial">Cashier Processing</text>
  <rect x="568" y="32" width="130" height="500" rx="6" fill="#f0fdf4" stroke="#bbf7d0" stroke-width="1"/>
  <text x="633" y="52" text-anchor="middle" font-size="9" font-weight="bold" fill="#15803d" font-family="Arial">STAGE 5</text>
  <text x="633" y="64" text-anchor="middle" font-size="9" fill="#16a34a" font-family="Arial">Handover</text>
  <rect x="708" y="32" width="104" height="500" rx="6" fill="#fef2f2" stroke="#fecaca" stroke-width="1"/>
  <text x="760" y="52" text-anchor="middle" font-size="9" font-weight="bold" fill="#b91c1c" font-family="Arial">SYSTEM</text>
  <!-- Stage 1 steps -->
  <rect x="16" y="90" width="114" height="30" rx="5" fill="#dbeafe" stroke="#3b82f6" stroke-width="1.5"/>
  <text x="73" y="109" text-anchor="middle" font-size="9" fill="#1e40af" font-family="Arial">Login / Register</text>
  <rect x="16" y="136" width="114" height="30" rx="5" fill="#dbeafe" stroke="#3b82f6" stroke-width="1.5"/>
  <text x="73" y="155" text-anchor="middle" font-size="9" fill="#1e40af" font-family="Arial">Browse Produk</text>
  <rect x="16" y="182" width="114" height="30" rx="5" fill="#dbeafe" stroke="#3b82f6" stroke-width="1.5"/>
  <text x="73" y="201" text-anchor="middle" font-size="9" fill="#1e40af" font-family="Arial">Tambah ke Cart</text>
  <rect x="16" y="228" width="114" height="30" rx="5" fill="#dbeafe" stroke="#3b82f6" stroke-width="1.5"/>
  <text x="73" y="247" text-anchor="middle" font-size="9" fill="#1e40af" font-family="Arial">Submit Checkout</text>
  <rect x="20" y="276" width="106" height="22" rx="11" fill="#e0e7ff" stroke="#818cf8" stroke-width="1.5"/>
  <text x="73" y="291" text-anchor="middle" font-size="8" font-weight="bold" fill="#4338ca" font-family="Arial">PENDING_APPROVAL</text>
  <line x1="73" y1="120" x2="73" y2="134" stroke="#6b7280" stroke-width="1.3" marker-end="url(#a2)"/>
  <line x1="73" y1="166" x2="73" y2="180" stroke="#6b7280" stroke-width="1.3" marker-end="url(#a2)"/>
  <line x1="73" y1="212" x2="73" y2="226" stroke="#6b7280" stroke-width="1.3" marker-end="url(#a2)"/>
  <line x1="73" y1="258" x2="73" y2="274" stroke="#3b82f6" stroke-width="1.5" marker-end="url(#a2)"/>
  <line x1="130" y1="105" x2="148" y2="105" stroke="#6b7280" stroke-width="1.3" stroke-dasharray="3" marker-end="url(#a2)"/>
  <!-- Stage 2 steps -->
  <rect x="156" y="90" width="114" height="30" rx="5" fill="#dcfce7" stroke="#22c55e" stroke-width="1.5"/>
  <text x="213" y="109" text-anchor="middle" font-size="9" fill="#15803d" font-family="Arial">Lihat Queue</text>
  <rect x="156" y="136" width="114" height="30" rx="5" fill="#dcfce7" stroke="#22c55e" stroke-width="1.5"/>
  <text x="213" y="155" text-anchor="middle" font-size="9" fill="#15803d" font-family="Arial">Review Order</text>
  <polygon points="213,196 248,216 213,236 178,216" fill="#dcfce7" stroke="#22c55e" stroke-width="1.5"/>
  <text x="213" y="212" text-anchor="middle" font-size="8" font-weight="bold" fill="#15803d" font-family="Arial">Approve?</text>
  <line x1="178" y1="216" x2="156" y2="216" stroke="#ef4444" stroke-width="1.3" marker-end="url(#a2r)"/>
  <rect x="16" y="204" width="100" height="24" rx="5" fill="#fee2e2" stroke="#ef4444" stroke-width="1.5"/>
  <text x="66" y="220" text-anchor="middle" font-size="9" fill="#b91c1c" font-family="Arial">CANCELLED</text>
  <text x="176" y="244" font-size="8" fill="#ef4444" font-family="Arial">Tidak</text>
  <line x1="248" y1="216" x2="288" y2="216" stroke="#16a34a" stroke-width="1.5" marker-end="url(#a2g)"/>
  <text x="262" y="211" font-size="8" fill="#16a34a" font-family="Arial">Ya</text>
  <rect x="156" y="264" width="114" height="42" rx="5" fill="#dcfce7" stroke="#22c55e" stroke-width="1.5"/>
  <text x="213" y="280" text-anchor="middle" font-size="8" fill="#15803d" font-family="Arial">Deduct Stok</text>
  <text x="213" y="292" text-anchor="middle" font-size="8" fill="#15803d" font-family="Arial">Start Timer 30min</text>
  <text x="213" y="304" text-anchor="middle" font-size="8" fill="#15803d" font-family="Arial">Generate QR</text>
  <line x1="213" y1="120" x2="213" y2="134" stroke="#6b7280" stroke-width="1.3" marker-end="url(#a2)"/>
  <line x1="213" y1="166" x2="213" y2="194" stroke="#6b7280" stroke-width="1.3" marker-end="url(#a2)"/>
  <line x1="213" y1="238" x2="213" y2="262" stroke="#16a34a" stroke-width="1.5" marker-end="url(#a2g)"/>
  <line x1="270" y1="282" x2="288" y2="282" stroke="#6b7280" stroke-width="1.3" stroke-dasharray="3" marker-end="url(#a2)"/>
  <!-- Stage 3: QR -->
  <rect x="296" y="90" width="114" height="30" rx="5" fill="#fef9c3" stroke="#ca8a04" stroke-width="1.5"/>
  <text x="353" y="109" text-anchor="middle" font-size="9" fill="#a16207" font-family="Arial">Terima WA Notif</text>
  <rect x="296" y="136" width="114" height="30" rx="5" fill="#dbeafe" stroke="#3b82f6" stroke-width="1.5"/>
  <text x="353" y="155" text-anchor="middle" font-size="9" fill="#1e40af" font-family="Arial">Buka Link QR</text>
  <rect x="296" y="182" width="114" height="30" rx="5" fill="#dbeafe" stroke="#3b82f6" stroke-width="1.5"/>
  <text x="353" y="201" text-anchor="middle" font-size="9" fill="#1e40af" font-family="Arial">Scan / Tunjuk QR</text>
  <rect x="300" y="228" width="106" height="22" rx="11" fill="#fef9c3" stroke="#ca8a04" stroke-width="1.5"/>
  <text x="353" y="243" text-anchor="middle" font-size="8" font-weight="bold" fill="#a16207" font-family="Arial">PENDING (Timer 30min)</text>
  <line x1="353" y1="120" x2="353" y2="134" stroke="#6b7280" stroke-width="1.3" marker-end="url(#a2)"/>
  <line x1="353" y1="166" x2="353" y2="180" stroke="#6b7280" stroke-width="1.3" marker-end="url(#a2)"/>
  <line x1="353" y1="212" x2="353" y2="226" stroke="#ca8a04" stroke-width="1.5" marker-end="url(#a2)"/>
  <line x1="410" y1="197" x2="428" y2="197" stroke="#6b7280" stroke-width="1.3" stroke-dasharray="3" marker-end="url(#a2)"/>
  <!-- Stage 4: Cashier -->
  <rect x="436" y="90" width="114" height="30" rx="5" fill="#fef9c3" stroke="#ca8a04" stroke-width="1.5"/>
  <text x="493" y="109" text-anchor="middle" font-size="9" fill="#a16207" font-family="Arial">Terima QR Scan</text>
  <rect x="436" y="136" width="114" height="30" rx="5" fill="#fef9c3" stroke="#ca8a04" stroke-width="1.5"/>
  <text x="493" y="155" text-anchor="middle" font-size="9" fill="#a16207" font-family="Arial">Verifikasi Detail</text>
  <rect x="436" y="182" width="114" height="30" rx="5" fill="#fef9c3" stroke="#ca8a04" stroke-width="1.5"/>
  <text x="493" y="201" text-anchor="middle" font-size="9" fill="#a16207" font-family="Arial">Proses Bayar</text>
  <rect x="440" y="228" width="106" height="22" rx="11" fill="#bbf7d0" stroke="#15803d" stroke-width="1.5"/>
  <text x="493" y="243" text-anchor="middle" font-size="8" font-weight="bold" fill="#14532d" font-family="Arial">PAID</text>
  <line x1="493" y1="120" x2="493" y2="134" stroke="#6b7280" stroke-width="1.3" marker-end="url(#a2)"/>
  <line x1="493" y1="166" x2="493" y2="180" stroke="#6b7280" stroke-width="1.3" marker-end="url(#a2)"/>
  <line x1="493" y1="212" x2="493" y2="226" stroke="#16a34a" stroke-width="1.5" marker-end="url(#a2g)"/>
  <line x1="550" y1="105" x2="568" y2="105" stroke="#6b7280" stroke-width="1.3" stroke-dasharray="3" marker-end="url(#a2)"/>
  <!-- Stage 5: Handover -->
  <rect x="576" y="90" width="114" height="30" rx="5" fill="#dcfce7" stroke="#22c55e" stroke-width="1.5"/>
  <text x="633" y="109" text-anchor="middle" font-size="9" fill="#15803d" font-family="Arial">Siapkan Barang</text>
  <rect x="576" y="136" width="114" height="30" rx="5" fill="#dcfce7" stroke="#22c55e" stroke-width="1.5"/>
  <text x="633" y="155" text-anchor="middle" font-size="9" fill="#15803d" font-family="Arial">Konfirmasi Handover</text>
  <rect x="576" y="192" width="114" height="22" rx="11" fill="#bfdbfe" stroke="#2563eb" stroke-width="1.5"/>
  <text x="633" y="207" text-anchor="middle" font-size="8" font-weight="bold" fill="#1e40af" font-family="Arial">COMPLETED</text>
  <line x1="633" y1="120" x2="633" y2="134" stroke="#6b7280" stroke-width="1.3" marker-end="url(#a2)"/>
  <line x1="633" y1="166" x2="633" y2="190" stroke="#2563eb" stroke-width="1.5" marker-end="url(#a2)"/>
  <!-- System column -->
  <rect x="716" y="90" width="88" height="30" rx="5" fill="#fee2e2" stroke="#ef4444" stroke-width="1.5"/>
  <text x="760" y="109" text-anchor="middle" font-size="9" fill="#b91c1c" font-family="Arial">WA Notif Approve</text>
  <rect x="716" y="136" width="88" height="30" rx="5" fill="#fee2e2" stroke="#ef4444" stroke-width="1.5"/>
  <text x="760" y="155" text-anchor="middle" font-size="9" fill="#b91c1c" font-family="Arial">WS Broadcast</text>
  <rect x="716" y="182" width="88" height="30" rx="5" fill="#fee2e2" stroke="#ef4444" stroke-width="1.5"/>
  <text x="760" y="201" text-anchor="middle" font-size="9" fill="#b91c1c" font-family="Arial">Auto Expired</text>
  <text x="760" y="224" text-anchor="middle" font-size="8" fill="#6b7280" font-family="Arial">jika 30min tidak bayar</text>
  <!-- Multi-booth note -->
  <rect x="148" y="368" width="272" height="60" rx="6" fill="#f5f3ff" stroke="#c4b5fd" stroke-width="1.5"/>
  <text x="284" y="385" text-anchor="middle" font-size="9" font-weight="bold" fill="#6d28d9" font-family="Arial">Multi-Booth Support</text>
  <text x="284" y="400" text-anchor="middle" font-size="8" fill="#6d28d9" font-family="Arial">Setiap helper booth approve secara independen</text>
  <text x="284" y="418" text-anchor="middle" font-size="8" fill="#6d28d9" font-family="Arial">Status PENDING hanya jika SEMUA booth approve</text>
</svg>"""

SVG_HELPER_INPUT = """<?xml version="1.0" encoding="UTF-8"?>
<svg viewBox="0 0 820 440" xmlns="http://www.w3.org/2000/svg">
  <rect width="820" height="440" fill="white"/>
  <defs>
    <marker id="a3" markerWidth="8" markerHeight="6" refX="8" refY="3" orient="auto"><polygon points="0 0,8 3,0 6" fill="#6b7280"/></marker>
    <marker id="a3g" markerWidth="8" markerHeight="6" refX="8" refY="3" orient="auto"><polygon points="0 0,8 3,0 6" fill="#16a34a"/></marker>
  </defs>
  <rect x="0" y="0" width="820" height="26" fill="#15803d"/>
  <text x="410" y="18" text-anchor="middle" font-size="12" font-weight="bold" fill="white" font-family="Arial">MODE: HELPER_INPUT — Booth Only (Staff-Assisted)</text>
  <rect x="8" y="32" width="186" height="400" rx="6" fill="#f0fdf4" stroke="#bbf7d0"/>
  <text x="101" y="48" text-anchor="middle" font-size="9" font-weight="bold" fill="#15803d" font-family="Arial">HELPER (Booth Staff)</text>
  <rect x="202" y="32" width="186" height="400" rx="6" fill="#fffbeb" stroke="#fde68a"/>
  <text x="295" y="48" text-anchor="middle" font-size="9" font-weight="bold" fill="#a16207" font-family="Arial">SYSTEM / AUTO</text>
  <rect x="396" y="32" width="186" height="400" rx="6" fill="#dbeafe" stroke="#bfdbfe"/>
  <text x="489" y="48" text-anchor="middle" font-size="9" font-weight="bold" fill="#1e40af" font-family="Arial">CUSTOMER</text>
  <rect x="590" y="32" width="222" height="400" rx="6" fill="#fff7ed" stroke="#fed7aa"/>
  <text x="701" y="48" text-anchor="middle" font-size="9" font-weight="bold" fill="#c2410c" font-family="Arial">CASHIER</text>
  <!-- Helper steps -->
  <rect x="16" y="60" width="170" height="28" rx="5" fill="#dcfce7" stroke="#22c55e" stroke-width="1.5"/>
  <text x="101" y="78" text-anchor="middle" font-size="9" fill="#15803d" font-family="Arial">Login Helper (Tablet)</text>
  <rect x="16" y="104" width="170" height="28" rx="5" fill="#dcfce7" stroke="#22c55e" stroke-width="1.5"/>
  <text x="101" y="122" text-anchor="middle" font-size="9" fill="#15803d" font-family="Arial">Cari / Input Customer</text>
  <rect x="16" y="148" width="170" height="28" rx="5" fill="#dcfce7" stroke="#22c55e" stroke-width="1.5"/>
  <text x="101" y="166" text-anchor="middle" font-size="9" fill="#15803d" font-family="Arial">Pilih &amp; Scan Produk</text>
  <rect x="16" y="192" width="170" height="28" rx="5" fill="#dcfce7" stroke="#22c55e" stroke-width="1.5"/>
  <text x="101" y="210" text-anchor="middle" font-size="9" fill="#15803d" font-family="Arial">Submit Order</text>
  <rect x="20" y="238" width="162" height="22" rx="11" fill="#dcfce7" stroke="#15803d" stroke-width="1.5"/>
  <text x="101" y="253" text-anchor="middle" font-size="8" font-weight="bold" fill="#15803d" font-family="Arial">RESERVED</text>
  <rect x="16" y="290" width="170" height="28" rx="5" fill="#dcfce7" stroke="#22c55e" stroke-width="1.5"/>
  <text x="101" y="308" text-anchor="middle" font-size="9" fill="#15803d" font-family="Arial">Konfirmasi Handover</text>
  <rect x="20" y="334" width="162" height="22" rx="11" fill="#bfdbfe" stroke="#2563eb" stroke-width="1.5"/>
  <text x="101" y="349" text-anchor="middle" font-size="8" font-weight="bold" fill="#1e40af" font-family="Arial">COMPLETED</text>
  <line x1="101" y1="88" x2="101" y2="102" stroke="#6b7280" stroke-width="1.3" marker-end="url(#a3)"/>
  <line x1="101" y1="132" x2="101" y2="146" stroke="#6b7280" stroke-width="1.3" marker-end="url(#a3)"/>
  <line x1="101" y1="176" x2="101" y2="190" stroke="#6b7280" stroke-width="1.3" marker-end="url(#a3)"/>
  <line x1="101" y1="220" x2="101" y2="236" stroke="#15803d" stroke-width="1.5" marker-end="url(#a3g)"/>
  <line x1="101" y1="318" x2="101" y2="332" stroke="#15803d" stroke-width="1.5" marker-end="url(#a3g)"/>
  <!-- System steps -->
  <rect x="210" y="60" width="170" height="28" rx="5" fill="#fef9c3" stroke="#ca8a04" stroke-width="1.5"/>
  <text x="295" y="78" text-anchor="middle" font-size="9" fill="#a16207" font-family="Arial">Validasi Stok</text>
  <rect x="210" y="104" width="170" height="28" rx="5" fill="#fef9c3" stroke="#ca8a04" stroke-width="1.5"/>
  <text x="295" y="122" text-anchor="middle" font-size="9" fill="#a16207" font-family="Arial">Potong Stok Langsung</text>
  <rect x="210" y="148" width="170" height="28" rx="5" fill="#fef9c3" stroke="#ca8a04" stroke-width="1.5"/>
  <text x="295" y="166" text-anchor="middle" font-size="9" fill="#a16207" font-family="Arial">Generate Token + QR</text>
  <rect x="210" y="192" width="170" height="28" rx="5" fill="#fef9c3" stroke="#ca8a04" stroke-width="1.5"/>
  <text x="295" y="210" text-anchor="middle" font-size="9" fill="#a16207" font-family="Arial">Kirim WA ke Customer</text>
  <line x1="295" y1="88" x2="295" y2="102" stroke="#6b7280" stroke-width="1.3" marker-end="url(#a3)"/>
  <line x1="295" y1="132" x2="295" y2="146" stroke="#6b7280" stroke-width="1.3" marker-end="url(#a3)"/>
  <line x1="295" y1="176" x2="295" y2="190" stroke="#6b7280" stroke-width="1.3" marker-end="url(#a3)"/>
  <line x1="186" y1="75" x2="208" y2="75" stroke="#6b7280" stroke-width="1.3" stroke-dasharray="3" marker-end="url(#a3)"/>
  <!-- Customer steps -->
  <rect x="404" y="104" width="170" height="28" rx="5" fill="#dbeafe" stroke="#3b82f6" stroke-width="1.5"/>
  <text x="489" y="122" text-anchor="middle" font-size="9" fill="#1e40af" font-family="Arial">Terima WA + Link QR</text>
  <rect x="404" y="148" width="170" height="28" rx="5" fill="#dbeafe" stroke="#3b82f6" stroke-width="1.5"/>
  <text x="489" y="166" text-anchor="middle" font-size="9" fill="#1e40af" font-family="Arial">Buka &amp; Lihat QR</text>
  <rect x="404" y="192" width="170" height="28" rx="5" fill="#dbeafe" stroke="#3b82f6" stroke-width="1.5"/>
  <text x="489" y="210" text-anchor="middle" font-size="9" fill="#1e40af" font-family="Arial">Pergi ke Kasir &amp; Bayar</text>
  <line x1="489" y1="132" x2="489" y2="146" stroke="#6b7280" stroke-width="1.3" marker-end="url(#a3)"/>
  <line x1="489" y1="176" x2="489" y2="190" stroke="#6b7280" stroke-width="1.3" marker-end="url(#a3)"/>
  <line x1="380" y1="206" x2="402" y2="118" stroke="#6b7280" stroke-width="1.3" stroke-dasharray="3" marker-end="url(#a3)"/>
  <!-- Cashier steps -->
  <rect x="598" y="104" width="198" height="28" rx="5" fill="#fff7ed" stroke="#f97316" stroke-width="1.5"/>
  <text x="697" y="122" text-anchor="middle" font-size="9" fill="#c2410c" font-family="Arial">Terima Customer + QR</text>
  <rect x="598" y="148" width="198" height="28" rx="5" fill="#fff7ed" stroke="#f97316" stroke-width="1.5"/>
  <text x="697" y="166" text-anchor="middle" font-size="9" fill="#c2410c" font-family="Arial">Proses Pembayaran</text>
  <rect x="598" y="192" width="198" height="28" rx="5" fill="#fff7ed" stroke="#f97316" stroke-width="1.5"/>
  <text x="697" y="210" text-anchor="middle" font-size="9" fill="#c2410c" font-family="Arial">Update Status PAID</text>
  <rect x="602" y="238" width="190" height="22" rx="11" fill="#bbf7d0" stroke="#15803d" stroke-width="1.5"/>
  <text x="697" y="253" text-anchor="middle" font-size="8" font-weight="bold" fill="#14532d" font-family="Arial">PAID → Helper Handover</text>
  <line x1="697" y1="132" x2="697" y2="146" stroke="#6b7280" stroke-width="1.3" marker-end="url(#a3)"/>
  <line x1="697" y1="176" x2="697" y2="190" stroke="#6b7280" stroke-width="1.3" marker-end="url(#a3)"/>
  <line x1="697" y1="220" x2="697" y2="236" stroke="#f97316" stroke-width="1.5" marker-end="url(#a3)"/>
  <line x1="574" y1="206" x2="596" y2="118" stroke="#6b7280" stroke-width="1.3" stroke-dasharray="3" marker-end="url(#a3)"/>
  <line x1="697" y1="260" x2="697" y2="290" stroke="#6b7280" stroke-width="1.3" stroke-dasharray="3"/>
  <line x1="697" y1="290" x2="186" y2="305" stroke="#6b7280" stroke-width="1.3" stroke-dasharray="3" marker-end="url(#a3)"/>
</svg>"""

SVG_SELF_ORDER = """<?xml version="1.0" encoding="UTF-8"?>
<svg viewBox="0 0 820 440" xmlns="http://www.w3.org/2000/svg">
  <rect width="820" height="440" fill="white"/>
  <defs>
    <marker id="a4" markerWidth="8" markerHeight="6" refX="8" refY="3" orient="auto"><polygon points="0 0,8 3,0 6" fill="#6b7280"/></marker>
    <marker id="a4g" markerWidth="8" markerHeight="6" refX="8" refY="3" orient="auto"><polygon points="0 0,8 3,0 6" fill="#16a34a"/></marker>
  </defs>
  <rect x="0" y="0" width="820" height="26" fill="#1d4ed8"/>
  <text x="410" y="18" text-anchor="middle" font-size="12" font-weight="bold" fill="white" font-family="Arial">MODE: SELF_ORDER — Kiosk / PWA (Customer-Driven)</text>
  <rect x="8" y="32" width="256" height="400" rx="6" fill="#dbeafe" opacity="0.4" stroke="#bfdbfe"/>
  <rect x="272" y="32" width="256" height="400" rx="6" fill="#fef9c3" opacity="0.4" stroke="#fde68a"/>
  <rect x="536" y="32" width="276" height="400" rx="6" fill="#f0fdf4" opacity="0.4" stroke="#bbf7d0"/>
  <text x="136" y="48" text-anchor="middle" font-size="9" font-weight="bold" fill="#1e40af" font-family="Arial">CUSTOMER</text>
  <text x="400" y="48" text-anchor="middle" font-size="9" font-weight="bold" fill="#a16207" font-family="Arial">SYSTEM AUTO</text>
  <text x="674" y="48" text-anchor="middle" font-size="9" font-weight="bold" fill="#15803d" font-family="Arial">CASHIER + HELPER</text>
  <!-- Customer column -->
  <rect x="16" y="60" width="240" height="28" rx="5" fill="#dbeafe" stroke="#3b82f6" stroke-width="1.5"/>
  <text x="136" y="78" text-anchor="middle" font-size="9" fill="#1e40af" font-family="Arial">Registrasi / Login (Phone)</text>
  <rect x="16" y="104" width="240" height="28" rx="5" fill="#dbeafe" stroke="#3b82f6" stroke-width="1.5"/>
  <text x="136" y="122" text-anchor="middle" font-size="9" fill="#1e40af" font-family="Arial">Browse Katalog Produk</text>
  <rect x="16" y="148" width="240" height="28" rx="5" fill="#dbeafe" stroke="#3b82f6" stroke-width="1.5"/>
  <text x="136" y="166" text-anchor="middle" font-size="9" fill="#1e40af" font-family="Arial">Add to Cart (max 20 item)</text>
  <rect x="16" y="192" width="240" height="28" rx="5" fill="#dbeafe" stroke="#3b82f6" stroke-width="1.5"/>
  <text x="136" y="210" text-anchor="middle" font-size="9" fill="#1e40af" font-family="Arial">Apply Voucher (Optional)</text>
  <rect x="16" y="236" width="240" height="28" rx="5" fill="#dbeafe" stroke="#3b82f6" stroke-width="1.5"/>
  <text x="136" y="254" text-anchor="middle" font-size="9" fill="#1e40af" font-family="Arial">Checkout &amp; Konfirmasi</text>
  <rect x="20" y="282" width="232" height="22" rx="11" fill="#fef9c3" stroke="#ca8a04" stroke-width="1.5"/>
  <text x="136" y="297" text-anchor="middle" font-size="8" font-weight="bold" fill="#a16207" font-family="Arial">PENDING (Timer: 8 menit)</text>
  <rect x="16" y="320" width="240" height="28" rx="5" fill="#dbeafe" stroke="#3b82f6" stroke-width="1.5"/>
  <text x="136" y="338" text-anchor="middle" font-size="9" fill="#1e40af" font-family="Arial">Scan QR / Bayar di Kasir</text>
  <line x1="136" y1="88" x2="136" y2="102" stroke="#6b7280" stroke-width="1.3" marker-end="url(#a4)"/>
  <line x1="136" y1="132" x2="136" y2="146" stroke="#6b7280" stroke-width="1.3" marker-end="url(#a4)"/>
  <line x1="136" y1="176" x2="136" y2="190" stroke="#6b7280" stroke-width="1.3" marker-end="url(#a4)"/>
  <line x1="136" y1="220" x2="136" y2="234" stroke="#6b7280" stroke-width="1.3" marker-end="url(#a4)"/>
  <line x1="136" y1="264" x2="136" y2="280" stroke="#ca8a04" stroke-width="1.5" marker-end="url(#a4)"/>
  <line x1="136" y1="304" x2="136" y2="318" stroke="#6b7280" stroke-width="1.3" marker-end="url(#a4)"/>
  <!-- System column -->
  <rect x="280" y="60" width="240" height="28" rx="5" fill="#fef9c3" stroke="#ca8a04" stroke-width="1.5"/>
  <text x="400" y="78" text-anchor="middle" font-size="9" fill="#a16207" font-family="Arial">Validasi Cart (Stok + On-Hold)</text>
  <rect x="280" y="104" width="240" height="28" rx="5" fill="#fef9c3" stroke="#ca8a04" stroke-width="1.5"/>
  <text x="400" y="122" text-anchor="middle" font-size="9" fill="#a16207" font-family="Arial">Potong Stok Langsung</text>
  <rect x="280" y="148" width="240" height="28" rx="5" fill="#fef9c3" stroke="#ca8a04" stroke-width="1.5"/>
  <text x="400" y="166" text-anchor="middle" font-size="9" fill="#a16207" font-family="Arial">Hitung PPN 12%</text>
  <rect x="280" y="192" width="240" height="28" rx="5" fill="#fef9c3" stroke="#ca8a04" stroke-width="1.5"/>
  <text x="400" y="210" text-anchor="middle" font-size="9" fill="#a16207" font-family="Arial">Generate QR + Timer 8 menit</text>
  <rect x="280" y="280" width="240" height="28" rx="5" fill="#fee2e2" stroke="#ef4444" stroke-width="1.5"/>
  <text x="400" y="298" text-anchor="middle" font-size="9" fill="#b91c1c" font-family="Arial">TxnExpireJob → EXPIRED</text>
  <text x="400" y="318" text-anchor="middle" font-size="8" fill="#b91c1c" font-family="Arial">jika tidak bayar dalam 8 menit</text>
  <line x1="400" y1="88" x2="400" y2="102" stroke="#6b7280" stroke-width="1.3" marker-end="url(#a4)"/>
  <line x1="400" y1="132" x2="400" y2="146" stroke="#6b7280" stroke-width="1.3" marker-end="url(#a4)"/>
  <line x1="400" y1="176" x2="400" y2="190" stroke="#6b7280" stroke-width="1.3" marker-end="url(#a4)"/>
  <line x1="256" y1="75" x2="278" y2="75" stroke="#6b7280" stroke-width="1.3" stroke-dasharray="3" marker-end="url(#a4)"/>
  <line x1="280" y1="206" x2="256" y2="295" stroke="#ca8a04" stroke-width="1.3" stroke-dasharray="3" marker-end="url(#a4)"/>
  <!-- Cashier + Helper column -->
  <rect x="544" y="60" width="260" height="28" rx="5" fill="#fef9c3" stroke="#ca8a04" stroke-width="1.5"/>
  <text x="674" y="78" text-anchor="middle" font-size="9" fill="#a16207" font-family="Arial">Cashier: Terima Scan QR</text>
  <rect x="544" y="104" width="260" height="28" rx="5" fill="#fef9c3" stroke="#ca8a04" stroke-width="1.5"/>
  <text x="674" y="122" text-anchor="middle" font-size="9" fill="#a16207" font-family="Arial">Cashier: Proses Bayar</text>
  <rect x="548" y="148" width="252" height="22" rx="11" fill="#bbf7d0" stroke="#15803d" stroke-width="1.5"/>
  <text x="674" y="163" text-anchor="middle" font-size="8" font-weight="bold" fill="#14532d" font-family="Arial">PAID</text>
  <rect x="544" y="182" width="260" height="28" rx="5" fill="#dcfce7" stroke="#22c55e" stroke-width="1.5"/>
  <text x="674" y="200" text-anchor="middle" font-size="9" fill="#15803d" font-family="Arial">Helper: Siapkan &amp; Handover</text>
  <rect x="548" y="226" width="252" height="22" rx="11" fill="#bfdbfe" stroke="#2563eb" stroke-width="1.5"/>
  <text x="674" y="241" text-anchor="middle" font-size="8" font-weight="bold" fill="#1e40af" font-family="Arial">COMPLETED</text>
  <line x1="674" y1="88" x2="674" y2="102" stroke="#6b7280" stroke-width="1.3" marker-end="url(#a4)"/>
  <line x1="674" y1="132" x2="674" y2="146" stroke="#16a34a" stroke-width="1.5" marker-end="url(#a4g)"/>
  <line x1="674" y1="170" x2="674" y2="180" stroke="#6b7280" stroke-width="1.3" marker-end="url(#a4)"/>
  <line x1="674" y1="210" x2="674" y2="224" stroke="#2563eb" stroke-width="1.5" marker-end="url(#a4)"/>
  <line x1="256" y1="334" x2="542" y2="75" stroke="#6b7280" stroke-width="1.3" stroke-dasharray="3" marker-end="url(#a4)"/>
</svg>"""

# ─── PPTX GENERATION ──────────────────────────────────────────────────────────

SLIDES = [
    {
        "title": "Global Order Flow",
        "subtitle": "Semua Order Mode — Swim Lane View",
        "svg": SVG_GLOBAL,
        "header_color": (30, 58, 95),
        "scale": 1.6,
    },
    {
        "title": "Mode: HELPER_APPROVE",
        "subtitle": "Model D (Active Mode) — Customer checkout → Helper approve → Cashier",
        "svg": SVG_HELPER_APPROVE,
        "header_color": (109, 40, 217),
        "scale": 1.6,
    },
    {
        "title": "Mode: HELPER_INPUT",
        "subtitle": "Booth Only (Staff-Assisted) — Helper input langsung untuk customer",
        "svg": SVG_HELPER_INPUT,
        "header_color": (21, 128, 61),
        "scale": 1.6,
    },
    {
        "title": "Mode: SELF_ORDER",
        "subtitle": "Kiosk / PWA (Customer-Driven) — Customer order mandiri via HP",
        "svg": SVG_SELF_ORDER,
        "header_color": (29, 78, 216),
        "scale": 1.6,
    },
]

COVER_SLIDE = {
    "title": "Amazing Toys Fair 2026",
    "subtitle": "Process Flow Diagrams",
    "details": [
        "Sistem Hybrid Order — 3 Mode × 7 Role",
        "",
        "Isi dokumen:",
        "  1.  Global Flow (Swim Lane — Semua Mode)",
        "  2.  Mode HELPER_APPROVE (Model D — Active)",
        "  3.  Mode HELPER_INPUT (Booth Only)",
        "  4.  Mode SELF_ORDER (Kiosk / PWA)",
        "",
        "Status Machine:",
        "  PENDING_APPROVAL → PENDING / RESERVED → PAID → HANDED_OVER → COMPLETED",
        "  Terminal: CANCELLED, EXPIRED",
        "",
        "CR History: CR-035 s/d CR-040",
    ],
}

def svg_to_png_bytes(svg_str: str, scale: float = 1.6) -> bytes:
    """Convert SVG string to PNG bytes using svglib + reportlab."""
    with tempfile.NamedTemporaryFile(suffix=".svg", delete=False, mode="w", encoding="utf-8") as f:
        f.write(svg_str)
        tmp_path = f.name
    try:
        drawing = svg2rlg(tmp_path)
        if drawing is None:
            raise RuntimeError("svg2rlg returned None")
        drawing.width *= scale
        drawing.height *= scale
        drawing.transform = (scale, 0, 0, scale, 0, 0)
        buf = io.BytesIO()
        renderPM.drawToFile(drawing, buf, fmt="PNG", bg=0xFFFFFF)
        return buf.getvalue()
    finally:
        os.unlink(tmp_path)


def add_cover_slide(prs: Presentation):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)

    # Dark background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42)

    W = prs.slide_width
    H = prs.slide_height

    # Title
    txb = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), W - Inches(1.2), Inches(1.2))
    tf = txb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = COVER_SLIDE["title"]
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)

    # Subtitle
    txb2 = slide.shapes.add_textbox(Inches(0.6), Inches(2.4), W - Inches(1.2), Inches(0.7))
    tf2 = txb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = COVER_SLIDE["subtitle"]
    r2.font.size = Pt(22)
    r2.font.color.rgb = RGBColor(148, 163, 184)

    # Details box
    txb3 = slide.shapes.add_textbox(Inches(1.4), Inches(3.3), W - Inches(2.8), H - Inches(4.0))
    tf3 = txb3.text_frame
    tf3.word_wrap = True
    for i, line in enumerate(COVER_SLIDE["details"]):
        if i == 0:
            p3 = tf3.paragraphs[0]
        else:
            p3 = tf3.add_paragraph()
        r3 = p3.add_run()
        r3.text = line
        r3.font.size = Pt(13)
        r3.font.color.rgb = RGBColor(203, 213, 225)

    # Footer
    txb4 = slide.shapes.add_textbox(Inches(0.4), H - Inches(0.6), W - Inches(0.8), Inches(0.4))
    tf4 = txb4.text_frame
    p4 = tf4.paragraphs[0]
    p4.alignment = PP_ALIGN.CENTER
    r4 = p4.add_run()
    r4.text = "Confidential — Internal Use Only"
    r4.font.size = Pt(9)
    r4.font.color.rgb = RGBColor(71, 85, 105)


def add_diagram_slide(prs: Presentation, slide_data: dict, slide_num: int):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)

    W = prs.slide_width
    H = prs.slide_height

    # Header bar
    hdr = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, W, Inches(0.72)
    )
    hdr.fill.solid()
    r, g, b = slide_data["header_color"]
    hdr.fill.fore_color.rgb = RGBColor(r, g, b)
    hdr.line.fill.background()

    # Slide number badge
    badge = slide.shapes.add_shape(1, Inches(0.2), Inches(0.1), Inches(0.44), Inches(0.44))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(255, 255, 255)
    badge.line.fill.background()
    txb_badge = slide.shapes.add_textbox(Inches(0.2), Inches(0.1), Inches(0.44), Inches(0.44))
    p_b = txb_badge.text_frame.paragraphs[0]
    p_b.alignment = PP_ALIGN.CENTER
    r_b = p_b.add_run()
    r_b.text = str(slide_num)
    r_b.font.size = Pt(13)
    r_b.font.bold = True
    r_b.font.color.rgb = RGBColor(r, g, b)

    # Header title
    txb_h = slide.shapes.add_textbox(Inches(0.8), Inches(0.06), W - Inches(1.2), Inches(0.36))
    p_h = txb_h.text_frame.paragraphs[0]
    r_h = p_h.add_run()
    r_h.text = slide_data["title"]
    r_h.font.size = Pt(18)
    r_h.font.bold = True
    r_h.font.color.rgb = RGBColor(255, 255, 255)

    # Subtitle
    txb_s = slide.shapes.add_textbox(Inches(0.8), Inches(0.42), W - Inches(1.2), Inches(0.26))
    p_s = txb_s.text_frame.paragraphs[0]
    r_s = p_s.add_run()
    r_s.text = slide_data["subtitle"]
    r_s.font.size = Pt(9)
    r_s.font.color.rgb = RGBColor(200, 220, 255)

    # Convert SVG → PNG and embed
    print(f"  Rendering: {slide_data['title']}...")
    png_bytes = svg_to_png_bytes(slide_data["svg"], scale=slide_data.get("scale", 1.6))
    img_stream = io.BytesIO(png_bytes)

    img_top = Inches(0.82)
    img_h = H - img_top - Inches(0.18)
    img_w = W - Inches(0.3)
    slide.shapes.add_picture(img_stream, Inches(0.15), img_top, width=img_w, height=img_h)

    # Page number footer
    txb_pg = slide.shapes.add_textbox(W - Inches(1.0), H - Inches(0.35), Inches(0.8), Inches(0.28))
    p_pg = txb_pg.text_frame.paragraphs[0]
    p_pg.alignment = PP_ALIGN.RIGHT
    r_pg = p_pg.add_run()
    r_pg.text = f"{slide_num} / {len(SLIDES)}"
    r_pg.font.size = Pt(8)
    r_pg.font.color.rgb = RGBColor(150, 150, 150)


def main():
    print("Generating ATF2026_Flow_Charts.pptx ...")

    prs = Presentation()
    prs.slide_width = Inches(13.33)   # widescreen 16:9
    prs.slide_height = Inches(7.5)

    print("  Adding cover slide...")
    add_cover_slide(prs)

    for i, slide_data in enumerate(SLIDES, start=1):
        add_diagram_slide(prs, slide_data, i)

    out_path = r"C:\Dev\amazing_toys_odoo_online - Hybrid\ATF2026_Flow_Charts.pptx"
    prs.save(out_path)
    print(f"\nDone! File saved to:\n  {out_path}")


if __name__ == "__main__":
    main()
